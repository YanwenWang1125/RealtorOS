from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme (RealtorOS branding)
PRIMARY_BLUE = RGBColor(37, 99, 235)  # #2563eb
SUCCESS_GREEN = RGBColor(16, 185, 129)  # #10b981
ACCENT_PURPLE = RGBColor(139, 92, 246)  # #8b5cf6
DARK_TEXT = RGBColor(30, 30, 30)
LIGHT_GRAY = RGBColor(243, 244, 246)

def add_title_slide(prs, title, subtitle=""):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_BLUE

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Subtitle
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        p = subtitle_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(28)
        p.font.color.rgb = RGBColor(255, 255, 255)

    return slide

def add_content_slide(prs, title, content_lines, notes=""):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_BLUE

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.5), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    for i, line in enumerate(content_lines):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = line
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_TEXT
        p.space_before = Pt(12)
        p.level = 0

    # Add notes
    if notes:
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = notes

    return slide

def add_section_header(prs, title, subtitle=""):
    """Add a section header slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background gradient effect (using rectangle)
    background = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0), Inches(10), Inches(7.5)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = ACCENT_PURPLE
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    p = title_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        p = subtitle_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(255, 255, 255)

    return slide

# Slide 1: Title Slide
add_title_slide(prs, "RealtorOS", "AI-Powered CRM for Real Estate Professionals\n\"Never Miss a Follow-Up Again\"")

# Slide 2: The Problem
add_content_slide(prs, "The Problem: Real Estate Agents Face Critical Challenges", [
    "📉 80% of leads are lost due to poor follow-up",
    "⏰ 5+ hours/week spent writing follow-up emails manually",
    "📧 No visibility into which clients are engaged",
    "🗓️ Missed appointments and forgotten tasks",
    "💰 Lost revenue from disorganized client management"
], notes="DESIGN NOTE: Add split-screen visual showing stressed agent vs missed opportunities/money flying away")

# Slide 3: The Solution
add_content_slide(prs, "Introducing RealtorOS", [
    "✅ Automated Follow-Ups: AI generates personalized emails",
    "✅ Real-Time Engagement: Know when clients open & click emails",
    "✅ Smart Task Management: Never forget a follow-up",
    "✅ Pipeline Visibility: Track every client from lead to closed",
    "✅ 80% Time Savings: Focus on selling, not admin work"
], notes="SCREENSHOT PLACEHOLDER: Add dashboard screenshot with highlighted features showing the main dashboard with KPIs visible")

# Slide 4: How It Works
add_content_slide(prs, "Simple 4-Step Process", [
    "1. Add Client → System auto-creates 5 follow-up tasks",
    "2. AI Writes Email → Personalized content in seconds",
    "3. Auto-Send → Scheduled delivery at optimal times",
    "4. Track Engagement → Real-time open & click notifications"
], notes="DESIGN NOTE: Add workflow diagram with icons showing the flow from client creation to engagement tracking")

# Slide 5: Architecture Overview
add_content_slide(prs, "Enterprise-Grade Architecture", [
    "Frontend: Next.js + TypeScript + Tailwind CSS",
    "Backend Services:",
    "  • Auth Service (Google OAuth)",
    "  • CRM Service (Client Management)",
    "  • Task Service (Automation)",
    "  • Email Service (AI + SendGrid)",
    "  • Webhook Service (Real-time Tracking)",
    "  • Analytics Service (Dashboards)",
    "Infrastructure: PostgreSQL, Redis, Docker",
    "AI Integration: OpenAI GPT-4 + SendGrid"
], notes="DESIGN NOTE: Add microservices architecture diagram showing all services and data flow between them")

# Slide 6: Feature #1 - Client Management
add_content_slide(prs, "Feature #1: Complete Client Relationship Management", [
    "📋 Full CRUD operations (Create, Read, Update, Delete)",
    "🎯 Pipeline stages: Lead → Negotiating → Under Contract → Closed",
    "🏠 Property tracking (address, type, price)",
    "🔍 Advanced search & filtering",
    "📝 Custom fields (flexible JSON storage)",
    "👥 Contact information management"
], notes="SCREENSHOT PLACEHOLDER: Add screenshots of client list view and client detail page showing the pipeline, custom fields, and property information")

# Slide 7: Feature #2 - AI Email Generation
add_content_slide(prs, "Feature #2: AI-Powered Email Automation", [
    "🤖 OpenAI GPT-4 Integration: Context-aware email generation",
    "✍️ Rich Text Editor: Full formatting control",
    "🎨 Custom Instructions: 'Mention the nearby park'",
    "🔄 Regenerate: Don't like it? Generate again",
    "📧 SendGrid Delivery: Professional, reliable sending",
    "🛡️ XSS Protection: Sanitized HTML content"
], notes="SCREENSHOT PLACEHOLDER: Add screenshot of email composer showing AI generation in action, with the rich text editor and custom instructions field visible. Show before/after of AI-generated email")

# Slide 8: Feature #3 - Engagement Tracking
add_content_slide(prs, "Feature #3: Real-Time Email Engagement Tracking", [
    "📊 Email Status Pipeline: Queued → Sent → Delivered → Opened → Clicked",
    "⚡ Real-Time Updates: 10-second polling for recent emails",
    "🔗 Click Tracking: Know which links clients engage with",
    "📅 Timestamp Tracking: Exact open/click times",
    "🔐 Webhook Security: ECDSA signature verification",
    "",
    "Value: 'Know exactly when to follow up'"
], notes="SCREENSHOT PLACEHOLDER: Add screenshot of email detail page with engagement timeline showing the status progression and timestamps")

# Slide 9: Feature #4 - Smart Task Management
add_content_slide(prs, "Feature #4: Automated Task Scheduling & Management", [
    "🤖 Auto-Creation: 5 tasks per new client",
    "  • Day 1 (High Priority)",
    "  • Day 3 (Medium Priority)",
    "  • Week 1, Week 2, Month 1 (Low Priority)",
    "📅 Calendar View: Drag-and-drop rescheduling",
    "📋 Table View: Sortable, filterable task list",
    "✅ Actions: Complete, Skip, Reschedule, Cancel",
    "🎨 Color-Coded: Priority-based visual indicators"
], notes="SCREENSHOT PLACEHOLDER: Add split screen showing calendar view on left and table view on right, highlighting the drag-and-drop functionality and color coding")

# Slide 10: Feature #5 - Dashboard & Analytics
add_content_slide(prs, "Feature #5: Real-Time Analytics & Insights", [
    "9 Key Metrics:",
    "  • Total Clients | Active Clients | Pending Tasks",
    "  • Completed Tasks | Emails Sent (Today/Week)",
    "  • Open Rate % | Click Rate % | Conversion Rate %",
    "",
    "Visualizations:",
    "  • 📊 Client Stage Distribution (Pie Chart)",
    "  • 📈 Email Engagement Metrics (Bar Chart)",
    "  • 📰 Activity Feed (Real-time updates)",
    "",
    "Auto-Refresh: Every 60 seconds"
], notes="SCREENSHOT PLACEHOLDER: Add screenshot of full dashboard showing all 9 KPI cards, pie chart, bar chart, and activity feed")

# Slide 11: Security & Compliance
add_content_slide(prs, "Enterprise-Grade Security", [
    "🔐 JWT Authentication: Secure token-based auth",
    "🔑 Google OAuth: Single sign-on integration",
    "🛡️ XSS Protection: DOMPurify HTML sanitization",
    "✍️ Webhook Signatures: ECDSA verification",
    "⏰ Timestamp Validation: 10-minute window",
    "🔒 HTTPS Only: Encrypted data transmission",
    "🗃️ Soft Delete: Data recovery capability",
    "👤 Multi-Tenancy Ready: Agent data isolation"
], notes="DESIGN NOTE: Add security icons with checkmarks for each feature")

# Slide 12: Technology Stack
add_content_slide(prs, "Built with Modern, Scalable Technologies", [
    "Frontend:",
    "  • Next.js 14 | TypeScript | React 18 | Tailwind CSS",
    "  • TanStack Query | Zustand | Zod | TipTap",
    "",
    "Backend:",
    "  • FastAPI | Python 3.11 | SQLAlchemy | Alembic",
    "  • Celery | Redis | OpenAI API | SendGrid",
    "",
    "Infrastructure:",
    "  • Docker | PostgreSQL 16 | Azure Container Apps",
    "",
    "Why This Matters: 'Production-ready, enterprise-scalable from day one'"
], notes="DESIGN NOTE: Add tech logos in organized sections for visual appeal")

# Slide 13: Live Demo
add_section_header(prs, "Live Demo", "See RealtorOS in Action")

# Slide 14: Demo Flow
add_content_slide(prs, "Demo Flow (10 minutes)", [
    "1. Create New Client (John Doe)",
    "2. View Auto-Generated Tasks (5 follow-ups)",
    "3. AI Email Generation (with custom instructions)",
    "4. Send & Track Email (real-time status updates)",
    "5. Calendar Management (drag-and-drop reschedule)",
    "6. Dashboard Analytics (live KPIs)"
], notes="DEMO NOTE: Have backup screenshots ready in case of live demo issues. Practice this flow multiple times before presentation.")

# Slide 15: ROI & Time Savings
add_content_slide(prs, "Quantifiable Results", [
    "Time Savings:",
    "  • Manual Process: 30 min/client × 20 clients = 10 hours/week",
    "  • With RealtorOS: 5 min/client × 20 clients = 1.7 hours/week",
    "  • Saved: 8.3 hours/week = 433 hours/year",
    "",
    "Revenue Impact:",
    "  • 8.3 hours/week freed up for selling",
    "  • Average agent: $75,000/year → $36/hour",
    "  • Value: $15,600/year per agent",
    "",
    "Engagement Boost:",
    "  • 80% improvement in follow-up consistency",
    "  • 35% higher client engagement rates"
], notes="DESIGN NOTE: Add before/after comparison chart showing the time savings visually")

# Slide 16: Current Features
add_content_slide(prs, "What's Available Today (MVP - v1.0)", [
    "✅ Fully Implemented:",
    "  • Client Management (CRUD + Pipeline)",
    "  • AI Email Generation (OpenAI GPT-4)",
    "  • Email Sending (SendGrid Integration)",
    "  • Real-Time Engagement Tracking",
    "  • Automated Task Creation (5 follow-ups)",
    "  • Calendar & Table Views",
    "  • Dashboard Analytics (9 KPIs)",
    "  • Google OAuth Login",
    "  • Microservices Architecture (7 services)",
    "  • Docker Containerization",
    "",
    "Status: 'Production-ready for pilot users'"
], notes="DESIGN NOTE: Use green checkmarks and organize in two columns for better readability")

# Slide 17: Roadmap Phase 2
add_content_slide(prs, "Roadmap - Phase 2 (Q1 2025)", [
    "Next 90 Days - Immediate Enhancements:",
    "",
    "📧 Email Templates: Pre-built templates library",
    "🚫 Unsubscribe Flow: One-click opt-out",
    "🎂 Birthday Reminders: Automated celebration emails",
    "🏠 Anniversary Tracking: Home purchase milestones",
    "🔄 Celery Automation: Full background task processing",
    "🔐 Multi-Tenancy: Complete agent data isolation",
    "📊 Advanced Analytics: Trend charts & forecasting",
    "📤 Bulk Operations: Import clients, mass email"
], notes="DESIGN NOTE: Add timeline with milestone markers showing the 90-day plan")

# Slide 18: Roadmap Phase 3
add_content_slide(prs, "Roadmap - Phase 3 (Q2-Q3 2025)", [
    "Future Vision - Game-Changing Features:",
    "",
    "👥 Team Collaboration: Multi-user/team support",
    "📱 SMS Integration: Text message follow-ups",
    "📆 Calendar Sync: Google Calendar integration",
    "📄 Document Management: Contract & listing storage",
    "📈 Advanced Reporting: Custom report builder",
    "🌐 Public API: Third-party integrations",
    "🔗 Zapier Integration: Connect 5,000+ apps"
], notes="DESIGN NOTE: Add feature cards with icons for each item")

# Slide 19: Roadmap Phase 4
add_content_slide(prs, "Roadmap - Phase 4 (Q4 2025)", [
    "Long-Term Vision - Industry Leadership:",
    "",
    "📱 Mobile Apps: iOS & Android native apps",
    "🗣️ Voice AI: Phone call transcription & analysis",
    "🤖 Chatbot: AI-powered client communication",
    "📊 Predictive Analytics: ML-based deal forecasting",
    "🌍 Multilingual: Support 10+ languages",
    "🎥 Video Messaging: Personalized video emails",
    "🏢 Enterprise: White-label solution for brokerages"
], notes="DESIGN NOTE: Add futuristic concept illustrations to show innovation")

# Slide 20: Competitive Advantages
add_content_slide(prs, "Why Choose RealtorOS?", [
    "Comparison vs Competitors:",
    "",
    "RealtorOS advantages:",
    "  ✅ AI Email Generation (GPT-4) - Only us!",
    "  ✅ Real-Time Tracking (Webhooks) - Others use batch",
    "  ✅ Microservices Architecture - Scalable from day 1",
    "  ✅ Modern Tech Stack (Next.js) - Others use legacy",
    "  ✅ Open Source Ready - Community-driven",
    "",
    "Pricing: $49/mo vs $69/mo (Follow Up Boss) or $1,000+/mo (BoomTown)"
], notes="DESIGN NOTE: Add comparison table showing RealtorOS vs Follow Up Boss, LionDesk, and BoomTown with checkmarks and X marks")

# Slide 21: Pricing Strategy
add_content_slide(prs, "Flexible Pricing for Every Agent", [
    "Starter - $49/month:",
    "  • 100 clients | 500 emails/month | Basic analytics | Email support",
    "",
    "Professional - $99/month:",
    "  • 500 clients | 2,000 emails/month | Advanced analytics",
    "  • Priority support | Custom templates",
    "",
    "Enterprise - Custom:",
    "  • Unlimited clients | Unlimited emails | White-label option",
    "  • Dedicated support | API access"
], notes="DESIGN NOTE: Add 3-tier pricing cards with visual hierarchy emphasizing the Professional tier")

# Slide 22: Target Market
add_content_slide(prs, "Who Is RealtorOS For?", [
    "Primary Markets:",
    "  🏡 Solo Agents: 1-10 clients/month",
    "  👥 Small Teams: 2-5 agents",
    "  📈 Growing Brokerages: 10-50 agents",
    "",
    "Secondary Markets:",
    "  🏢 Loan Officers: Similar follow-up needs",
    "  💼 Property Managers: Tenant communication",
    "  🏗️ Commercial Agents: Long sales cycles",
    "",
    "Market Size: 1.5M real estate agents in US"
], notes="DESIGN NOTE: Add agent personas with photos representing each target segment")

# Slide 23: Go-To-Market Strategy
add_content_slide(prs, "Launch Plan", [
    "Phase 1 - Beta (Month 1-2):",
    "  • 50 pilot users (local agents)",
    "  • Feedback collection | Bug fixes & polish",
    "",
    "Phase 2 - Public Launch (Month 3-4):",
    "  • Product Hunt launch | SEO-optimized website",
    "  • YouTube demos & tutorials | Facebook Ads (real estate groups)",
    "",
    "Phase 3 - Scale (Month 5-12):",
    "  • Partnership with NAR/realtor.com",
    "  • Integration marketplace | Affiliate program (20% commission)"
], notes="DESIGN NOTE: Add marketing funnel diagram showing the progression from beta to scale")

# Slide 24: Metrics & Traction
add_content_slide(prs, "Current Status & Projections", [
    "Current (MVP):",
    "  ✅ MVP Complete (4 modules)",
    "  ✅ 5,000+ lines of documentation",
    "  ✅ Microservices architecture deployed",
    "",
    "6-Month Goals:",
    "  • 100 paying customers | $5,000 MRR",
    "  • 4.5+ star rating | 10,000 emails sent/day",
    "",
    "12-Month Goals:",
    "  • 1,000 customers | $50,000 MRR",
    "  • Mobile app launch | 50,000 emails/day"
], notes="DESIGN NOTE: Add growth charts showing projections if applicable, or use aspirational growth curves")

# Slide 25: Call to Action
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = SUCCESS_GREEN

# Title
title_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "Ready to Transform Your Real Estate Business?"
p = title_frame.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)

# CTAs
cta_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(3))
text_frame = cta_box.text_frame
ctas = [
    "🚀 Join Beta: [Website URL]",
    "📧 Book Demo: [Email/Calendly Link]",
    "💬 Contact Us: [Phone/Email]",
    "📱 Follow: [Social Media Handles]",
    "",
    "Special Offer: First 100 users get 50% off for 6 months"
]

for i, line in enumerate(ctas):
    if i > 0:
        text_frame.add_paragraph()
    p = text_frame.paragraphs[i]
    p.text = line
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(20) if i < len(ctas) - 1 else Pt(24)
    p.font.bold = (i == len(ctas) - 1)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.space_before = Pt(8)

notes_slide = slide.notes_slide
text_frame = notes_slide.notes_text_frame
text_frame.text = "DESIGN NOTE: Add large, clear CTA buttons and a QR code linking to demo/signup page"

# Slide 26: Q&A
add_section_header(prs, "Questions & Discussion", "Thank you!")

# Add notes to Q&A slide
last_slide = prs.slides[-1]
notes_slide = last_slide.notes_slide
text_frame = notes_slide.notes_text_frame
text_frame.text = """Be Ready to Answer:
- How does AI personalization work?
- What happens if AI generates bad content?
- How do you ensure email deliverability?
- What's the data migration process?
- Can it integrate with my existing CRM?
- What's the learning curve?
- How do you handle GDPR/privacy?"""

# Save presentation
output_path = r"c:\Users\Leo\AI projects\RealtorOS\RealtorOS_Demo_Presentation.pptx"
prs.save(output_path)
print(f"Presentation created successfully: {output_path}")
print(f"Total slides: {len(prs.slides)}")
